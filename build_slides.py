#!/usr/bin/env python3
"""Generate the course slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image as PILImage   # to fit screenshots inside their card

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_topic1 import TOPIC1
from data_topic2 import TOPIC2
from data_topic3 import TOPIC3
from data_topic4 import TOPIC4
from data_topic5 import TOPIC5
ACTIVITIES = TOPIC1 + TOPIC2 + TOPIC3 + TOPIC4 + TOPIC5

REPO = HERE          # generators live in the course repo root
ASSETS = os.path.join(REPO, "courseware", "assets")

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — driven by course_data
    rect(s,Inches(11.0),Inches(0.72),Inches(1.55),Inches(1.0),BLUE)
    txt(s,Inches(11.0),Inches(0.82),Inches(1.55),Inches(0.5),[[(C.BADGE_TOP,20,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(11.0),Inches(1.28),Inches(1.55),Inches(0.4),[[(C.BADGE_SUB,8,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,21,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Tools & libraries:  ",13,GREY,True),(services,13,GREY,False)]]); footer(s); return s
CODE_BG=RGBColor(0x0B,0x12,0x20); CODE_FG=RGBColor(0xF8,0xFA,0xFC)
CODE_COMMENT=RGBColor(0x7E,0x9A,0xB8); CODE_PROMPT=RGBColor(0x5E,0xEA,0xD4)
MAX_CODE_LINES=13          # lines that fit one dark card at 12pt above the footer
WRAP=62                    # characters before a code line is soft-wrapped
                           # 62 ch x 12pt Consolas (0.1in/ch) = 6.2in, inside the
                           # 7.55in card minus its 0.5in horizontal padding

def _wrap_code(cmd):
    """Split a snippet into display lines, soft-wrapping over-long ones.

    A wrapped continuation keeps the original line's indentation plus a hanging
    indent, so Python block structure survives the wrap — a dict body must never
    come back at column 0 and read as top-level code.

    Every pass must consume at least one character of the remainder, otherwise a
    line with no break opportunity (long JSON, curl payloads) would loop forever.
    """
    out=[]
    for raw in cmd.split("\n"):
        if len(raw)<=WRAP:
            out.append(raw); continue
        indent=min(len(raw)-len(raw.lstrip()), 8)
        pad=" "*(indent+2)
        cur=raw; first=True
        while len(cur)>WRAP:
            # continuations carry `pad`, so they have less room for real content
            avail=WRAP if first else max(WRAP-len(pad), 24)
            # never break inside the leading indent of a continuation line
            lo=0 if first else len(pad)
            # Prefer a space; dense code (.agg(['count','mean','median'])) often has
            # none, so fall back to a punctuation boundary rather than hard-cutting
            # mid-token — splitting 'median' into 'medi'/'an' is unreadable.
            cut=cur.rfind(" ", lo, avail)
            if cut<=lo:
                cut=max(cur.rfind(c, lo, avail)+1 for c in ",;([{")
            if cut<=lo: cut=avail                # hard cut — guarantees progress
            out.append(cur[:cut])
            rest=cur[cut:].lstrip()
            if not rest: break
            cur=pad+rest; first=False
        else:
            out.append(cur)
    return out

CODE_TOP=3.05              # where the code card starts
CODE_BOTTOM_LIMIT=6.55     # code must never cross into the footer band

def code_card(s,x,y,w,lines,size=12):
    """Dark card sized to its content, white/light code text.

    The card is clamped so it can never run past CODE_BOTTOM_LIMIT: callers
    chunk to MAX_CODE_LINES, and this is the backstop that keeps a miscount
    from spilling code over the footer.
    """
    lh=0.235 if size>=12 else 0.205
    h=min(0.34+lh*len(lines), CODE_BOTTOM_LIMIT-(y.inches if hasattr(y,"inches") else CODE_TOP))
    rect(s,x,y,w,Inches(h),CODE_BG)
    tb=s.shapes.add_textbox(x+Inches(0.28),y+Inches(0.16),w-Inches(0.5),Inches(h-0.3))
    tf=tb.text_frame; tf.word_wrap=False
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(0); p.line_spacing=1.0
        st=ln.strip()
        col=CODE_COMMENT if st.startswith("#") else (CODE_PROMPT if st.startswith("$") else CODE_FG)
        r=p.add_run(); r.text=ln if ln.strip() else " "
        r.font.size=Pt(size); r.font.name="Consolas"; r.font.color.rgb=col
    return h

def screenshot_slide(title,img,caption,kicker,notes=()):
    """Portal screenshot on the left, guidance on the right.

    The house standard wants the LMS/exam portals shown as a real screenshot —
    a bare URL in a bullet is not acceptable for these admin slides.
    """
    s=head(slide(),title,kicker)
    p=_logo(img)
    CX,CY,CW_,CH_=0.85,2.05,7.55,3.95      # host card
    if p:
        rect(s,Inches(CX),Inches(CY),Inches(CW_),Inches(CH_),LIGHT)
        # fit INSIDE the card (never fit-to-width): a wide screenshot would
        # otherwise overhang the card bottom and collide with the caption
        iw,ih=PILImage.open(p).size
        pad=0.15
        bw,bh=CW_-2*pad, CH_-2*pad
        sc=min(bw/iw, bh/ih)
        dw,dh=iw*sc, ih*sc
        s.shapes.add_picture(p,Inches(CX+(CW_-dw)/2),Inches(CY+(CH_-dh)/2),
                             width=Inches(dw),height=Inches(dh))
    # caption sits BELOW the card, clear of the image
    txt(s,Inches(CX),Inches(CY+CH_+0.12),Inches(CW_),Inches(0.3),[[(caption,12,GREY,False)]])
    ex=Inches(8.65); rect(s,ex,Inches(2.05),Inches(3.85),Inches(4.25),LIGHT)
    tb=s.shapes.add_textbox(ex+Inches(0.28),Inches(2.3),Inches(3.3),Inches(3.8))
    tf=tb.text_frame; tf.word_wrap=True
    for i,n_ in enumerate(notes):
        p_=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p_.space_after=Pt(10)
        r=p_.add_run(); r.text="•  "+n_
        r.font.size=Pt(14); r.font.name="Arial"; r.font.color.rgb=INK
    footer(s); return s

def _explain(cmd, instr):
    """Plain-language notes about the snippet, shown beside the code card."""
    if not cmd:
        return []
    txt_l=cmd.lower(); out=[]
    add=lambda t: out.append(t) if t not in out else None
    if txt_l.lstrip().startswith("$") or any(
            txt_l.lstrip().startswith(p) for p in ("uv ","curl","docker","git ","grep","open ")):
        add("Run this in your terminal, from the project folder.")
    if "uv add" in txt_l:   add("uv resolves, installs and records the dependency in uv.lock.")
    if "uv run" in txt_l:   add("uv run executes inside the project environment — no venv activation.")
    if "def " in cmd:       add("Read the signature first: what goes in, what comes out.")
    if '"""' in cmd:        add("The docstring states the contract the function promises.")
    if "raise " in cmd:     add("Invalid input fails fast and loudly, never silently.")
    if "class " in cmd:     add("The class bundles data with the behaviour that acts on it.")
    if "@property" in cmd:  add("A property exposes state read-only — callers cannot assign to it.")
    if "@abstractmethod" in cmd: add("Subclasses MUST implement this — the contract is enforced.")
    if "try:" in cmd:       add("try/except handles the failure; finally always cleans up.")
    if "groupby" in txt_l:  add("split-apply-combine: group the rows, then aggregate each group.")
    if "basemodel" in txt_l:add("Pydantic validates the data at runtime from the type hints.")
    if "@app." in cmd:      add("FastAPI turns this typed function into an HTTP endpoint.")
    if "assert " in cmd:    add("The assertion is the test — it must fail before you trust it.")
    # ---- framework note: pick by SCORE, not by the order these lines happen to
    # sit in. `"st." in cmd` also matches inside `pytest.`, and an `httpx` install
    # line is not a test — deciding on the first rule to fire put a pytest caption
    # on a FastAPI install slide and a Streamlit caption on pytest code.
    # Count real occurrences per framework and let the dominant one speak.
    import re as _re
    FRAMEWORKS=[
        ("streamlit", r"\bstreamlit\b|(?<!py)(?<!\w)st\.\w",
         "Streamlit renders each widget top-to-bottom and re-runs the script on every interaction."),
        ("pytest",    r"\bpytest\b|\bdef test_|\bassert\b",
         "The test drives the API exactly as a real client would."),
        ("fastapi",   r"@app\.",
         "FastAPI turns this typed function into an HTTP endpoint."),
        ("docker",    r"\bdocker\b|\bdockerfile\b|\bcompose\b|^FROM \w|^RUN |^CMD ",
         "The image bundles interpreter, dependencies and code as one artifact."),
        ("pydantic",  r"\bbasemodel\b|\bpydantic\b",
         "Pydantic validates the data at runtime from the type hints."),
        ("pandas",    r"\bgroupby\b|\brolling\b|\bagg\(",
         "split-apply-combine: group the rows, then aggregate each group."),
    ]
    scored=[(len(_re.findall(pat, cmd, _re.I|_re.M)), note) for _,pat,note in FRAMEWORKS]
    best=max(scored, key=lambda t: t[0])
    if best[0]:
        # the dominant framework note leads; keep at most two generic notes after it
        out=[best[1]]+[o for o in out if o not in {n for _,_,n in FRAMEWORKS}][:2]
    # Fall back to the step's own instruction when nothing matched.
    if not out:             add(instr)
    return out[:5]

def step_slide(kicker,act_title,n,total,text,cmd=""):
    """Code card on the left, explanation on the right — never overlapping.
    Long snippets continue onto further slides."""
    lines=_wrap_code(cmd) if cmd else []
    chunks=[lines[i:i+MAX_CODE_LINES] for i in range(0,len(lines),MAX_CODE_LINES)] or [[]]
    for ci,chunk in enumerate(chunks):
        # Notes describe the code ON THIS SLIDE. Deriving them once from the whole
        # snippet put a "try/except handles the failure" caption on a continuation
        # slide whose chunk contains no try: at all.
        notes=_explain("\n".join(chunk) or cmd, text)
        s=head(slide(),act_title,kicker,TEAL)
        label=f"STEP {n} OF {total}"+(f"  ·  CONT. {ci+1}/{len(chunks)}" if len(chunks)>1 else "")
        oval(s,Inches(0.85),Inches(1.95),Inches(0.82),Inches(0.82),TEAL)
        txt(s,Inches(0.85),Inches(2.08),Inches(0.82),Inches(0.6),
            [[(str(n),22,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(1.85),Inches(1.92),Inches(10.6),Inches(0.3),[[(label,11,GREY,True)]])
        txt(s,Inches(1.85),Inches(2.20),Inches(10.6),Inches(0.62),
            [[(text,19,INK,False)]],anchor=MSO_ANCHOR.TOP)
        if not chunk:
            footer(s); continue
        CW=Inches(7.55)                       # code column
        code_card(s,Inches(0.85),Inches(3.05),CW,chunk,size=12)
        # explanation column
        ex=Inches(8.65); rect(s,ex,Inches(3.05),Inches(3.85),Inches(3.6),LIGHT)
        txt(s,ex+Inches(0.28),Inches(3.25),Inches(3.3),Inches(0.35),
            [[("WHAT THIS DOES",11,BLUE,True)]])
        tb=s.shapes.add_textbox(ex+Inches(0.28),Inches(3.62),Inches(3.3),Inches(2.9))
        tf=tb.text_frame; tf.word_wrap=True
        for i,note_ in enumerate(notes):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.space_after=Pt(9)
            r=p.add_run(); r.text="•  "+note_
            r.font.size=Pt(13); r.font.name="Arial"; r.font.color.rgb=INK
        footer(s)
    return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
flow_h("Digital Attendance (Mandatory)",[
 "Mandatory at AM, PM and Assessment for every WSQ-funded course",
 "The trainer displays the SSG digital attendance QR code",
 "Scan it with your mobile phone camera",
 "Submit — your attendance is recorded on the SSG portal",
 "75% attendance minimum to be eligible for assessment and funding"],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 C.TRAINER_BULLETS,
 initials="AA",accent=BLUE)
content("Let's Know Each Other",C.ICE_BREAKER,kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
screenshot_slide("Download Your Course Material","lms-tms-portal.png",
 "https://lms-tms.tertiaryinfotech.com",kicker="COURSE PORTAL",
 notes=["Sign in with the email address you used to register for this course.",
        "Download the slides and the Learner Guide before the assessment.",
        "The assessment is OPEN BOOK — have both open on the day.",
        "Your attendance and assessment results are recorded here too."])
# Lesson plan overview
def _lab_span(topic_num):
    nums=[a["num"] for a in ACTIVITIES if a["topic"]==topic_num]
    return f"Labs {min(nums)}–{max(nums)}" if len(nums)>1 else f"Lab {nums[0]}"
two_col(f"Lesson Plan — {C.DAYS} Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 (f"Topic 1: {C.TOPICS[0]['title']} ({_lab_span(1)})",1),
 (f"Topic 2: {C.TOPICS[1]['title']} ({_lab_span(2)})",1),
 (f"Day 2 — {C.DAY_THEMES[2]}",0),
 (f"Topic 3: {C.TOPICS[2]['title']} ({_lab_span(3)})",1),
 (f"Topic 4: {C.TOPICS[3]['title']} begins ({_lab_span(4)})",1)],
 [(f"Day 3 — {C.DAY_THEMES[3]}",0),
 (f"Topic 4: {C.TOPICS[3]['title']} completes",1),
 (f"Topic 5: {C.TOPICS[4]['title']} ({_lab_span(5)})",1),
 ("Final Assessment (WA + PP)",1),
 ("Daily timing",0),
 ("9:30am–6:30pm · 1-hour lunch · tea breaks within",1)],
 kicker="SCHEDULE",lhead="Days 1–2",rhead="Day 3 & timing")
tile_grid("Learning Outcomes",
 [(f"LO{i}", lo.split(": ",1)[1] if ": " in lo else lo)
  for i,lo in enumerate(C.LEARNING_OUTCOMES,1)],
 kicker="WHAT YOU'LL ACHIEVE",cols=1,size=13)
tile_grid("Briefing for Assessment",[
 ("Clear your desk","Phones and materials go under the table or on the floor."),
 ("No photos or recording","Assessment scripts must not be photographed or recorded."),
 ("No discussion","Work individually — no talking during the assessment."),
 ("Black or blue pen","Use a black/blue pen for hard-copy assessments."),
 ("No correction fluid","Liquid paper and correction tape are not permitted."),
 ("Scripts collected","All scripts are collected when time is up.")],
 kicker="BEFORE YOU START",cols=2,size=15)
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study (CS) — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
# No Practice Exam slide: this course has no practice exam on
# exams.tertiaryinfotech.com, so the house-standard slide is deliberately omitted
# rather than pointing learners at a page that does not exist.

# ---------------- CORE CONCEPTS ----------------
section("CORE CONCEPTS","What Vibe Coding Is — and What It Is Not","")
tile_grid("What is Vibe Coding?",[
 ("AI-assisted, human-owned","You specify, the assistant drafts, you review and own the result."),
 ("Specification over guesswork","State the goal, constraints, inputs and expected output."),
 ("Verify, never assume","Generated code that runs is not the same as generated code that is correct."),
 ("Faster iteration","Refactoring conversationally beats rewriting by hand — with a test to prove it.")],
 kicker="OVERVIEW",cols=2,size=15)
cards3("The Vibe Coding Loop",[
 (BLUE,"1 · Specify",["State the goal precisely","Name the constraints","Give real inputs and expected output"]),
 (TEAL,"2 · Generate",["Let the assistant draft","Read every line before running","Ask for the shape you want"]),
 (VIOLET,"3 · Verify",["Run it against edge cases","Test the error paths","Fix by naming the defect, not 'improve this'"])],kicker="THE METHOD")
big_statement("You are accountable for every line you ship.","The assistant drafts; you review, test and own the result. That is the difference between vibe coding and guessing.","WHY IT MATTERS",color=BLUE)
two_col("The Toolchain for This Course",[
 ("uv",0),("Project, interpreter and locked dependencies",1),("uv init · uv add · uv run",1),
 ("pandas 3.0",0),("DataFrames, cleaning, groupby, rolling windows",1),("Copy-on-Write is the default",1)],
 [("Pydantic + FastAPI",0),("Typed models, validation, HTTP endpoints",1),("OpenAPI docs for free",1),
 ("Streamlit + Docker",0),("Analyst UI and a deployable image",1),("One image, any machine",1)],
 kicker="TOOLING",lhead="Build",rhead="Serve & ship")
content("Your Workbench",[
 "uv manages Python, the virtual environment and the lockfile — no venv activation by hand.",
 "An AI coding assistant (Claude, Copilot or Cursor) drafts code you review.",
 "SQLite ships with Python — the database needs no server and runs offline.",
 "Every lab in this course builds one application: CardGuard, a card-transaction fraud screening system."],kicker="YOUR WORKBENCH")

# ---------------- TOPICS + ACTIVITIES ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    # concept slide(s) — visual tile grid instead of a bullet list
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker="KEY CONCEPTS", cols=2, size=14)
    acts=TOPIC_ACTS[t["num"]]
    # a card summary of the labs in this topic
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        cards.append((CARD_COLORS[gi], f"Labs {g[0]['num']}–{g[-1]['num']}" if g else "—",
                      [a["title"] for a in g] if g else ["—"]))
    cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    # per activity
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} · HANDS-ON")
        # steps may be ("instruction", "code") or a bare ("instruction",) / "instruction"
        steps=[(s[0], s[1] if len(s)>1 else "") if isinstance(s,(tuple,list)) else (s,"")
               for s in a["steps"]]
        total=len(steps)
        for i,(instr,cmd) in enumerate(steps,1):
            # kicker is the lab number only — the full title is the H1 directly
            # beneath it, so truncating it here just cut words mid-token
            step_slide(f"LAB {a['num']}", a["title"], i, total, instr, cmd)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    # topic recap
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",
 [(t["title"], t["subtitle"]) for t in C.TOPICS],
 kicker="LEARNING OUTCOMES",cols=1,size=13)
content(C.NEXT_STEPS_TITLE, C.NEXT_STEPS, kicker="NEXT STEPS")
content("Keep Practising",[
 "Re-run every lab from a clean folder — uv sync restores the exact environment from uv.lock.",
 "Read the Learner Guide alongside your own code; the step numbering matches the labs exactly.",
 "Practise the four-part prompt pattern: goal, constraints, inputs, expected output.",
 "Tertiary Infotech practice exams: https://exams.tertiaryinfotech.com/"],kicker="TEST YOURSELF")
content("Assessment",[
 "Written Assessment (SAQ) — 50 minutes.  Case Study (CS) — 80 minutes.",
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="WRAP-UP")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then the Case Study (CS) — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
flow_h("Digital Attendance (Mandatory)",[
 "Mandatory at AM, PM and Assessment for every WSQ-funded course",
 "The trainer displays the SSG digital attendance QR code",
 "Scan it with your mobile phone camera",
 "Submit — your attendance is recorded on the SSG portal",
 "75% attendance minimum to be eligible for assessment and funding"],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!",C.THANK_YOU_LINE,C.THANK_YOU_KICKER,color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
